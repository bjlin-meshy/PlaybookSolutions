#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将HTML演示文稿转换为PowerPoint PPTX文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re

def clean_text(text):
    """清理文本，移除HTML标签和多余空白"""
    if not text:
        return ""
    # 移除HTML标签
    text = re.sub(r'<[^>]+>', '', text)
    # 清理空白字符，但保留换行
    text = re.sub(r'[ \t]+', ' ', text)
    text = text.replace('\n', ' ').strip()
    return text

def extract_text_from_element(element, lang='cn'):
    """从元素中提取指定语言的文本"""
    if not element:
        return ""
    
    # 查找指定语言的span
    lang_spans = element.find_all('span', class_=lang)
    if lang_spans:
        texts = []
        for span in lang_spans:
            text = clean_text(span.get_text())
            if text:
                texts.append(text)
        if texts:
            return '\n'.join(texts) if len(texts) > 1 else texts[0]
    
    # 如果没有找到语言特定的span，返回所有文本
    text = clean_text(element.get_text())
    return text

def add_text_to_shape(shape, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """向形状添加文本"""
    if not text:
        return
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    # 清除默认段落
    text_frame.clear()
    
    # 处理多行文本
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = line.strip()
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        if color:
            p.font.color.rgb = RGBColor(*color)
        else:
            p.font.color.rgb = RGBColor(255, 255, 255)

def create_slide_from_html(prs, slide_element, lang='cn'):
    """从HTML创建PowerPoint幻灯片"""
    # slide_element 已经是 BeautifulSoup 元素
    slide_content = slide_element.find('div', class_='slide-content')
    
    if not slide_content:
        return
    
    # 创建新幻灯片（使用空白布局）
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色为深灰色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 8, 8)
    
    # 检查是否是封面页或章节标题页
    is_cover = 'align-items: center' in str(slide_content.get('style', ''))
    is_section_header = 'section-header-slide' in slide_element.get('class', [])
    
    if is_cover or is_section_header:
        # 居中布局
        # 获取标题
        title_elem = slide_content.find('h1') or slide_content.find('h2')
        title_text = ""
        if title_elem:
            title_text = extract_text_from_element(title_elem, lang)
        
        # 添加主标题
        if title_text:
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(2)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            add_text_to_shape(title_shape, title_text, font_size=48, bold=True, 
                            color=(255, 255, 255), alignment=PP_ALIGN.CENTER)
        
        # 获取描述
        desc_elem = slide_content.find('p', class_='slide-desc')
        if not desc_elem:
            desc_elem = slide_content.find('p')
        if desc_elem:
            desc_text = extract_text_from_element(desc_elem, lang)
            if desc_text:
                left = Inches(1)
                top = Inches(4.5)
                width = Inches(8)
                height = Inches(1)
                desc_shape = slide.shapes.add_textbox(left, top, width, height)
                add_text_to_shape(desc_shape, desc_text, font_size=20, 
                                color=(136, 136, 136), alignment=PP_ALIGN.CENTER)
    else:
        # 常规内容页
        y_pos = 0.5
        
        # 获取section tag
        section_tag = slide_content.find('div', class_='section-tag')
        if section_tag:
            section_text = extract_text_from_element(section_tag, lang)
            if section_text:
                left = Inches(0.5)
                top = Inches(y_pos)
                width = Inches(9)
                height = Inches(0.4)
                tag_shape = slide.shapes.add_textbox(left, top, width, height)
                add_text_to_shape(tag_shape, section_text, font_size=12, bold=True, 
                                color=(204, 255, 0))
                y_pos += 0.5
        
        # 获取标题
        title_elem = slide_content.find('h1') or slide_content.find('h2', class_='slide-title')
        if not title_elem:
            title_elem = slide_content.find('h2')
        title_text = ""
        if title_elem:
            title_text = extract_text_from_element(title_elem, lang)
        
        if title_text:
            left = Inches(0.5)
            top = Inches(y_pos)
            width = Inches(9)
            height = Inches(0.8)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            add_text_to_shape(title_shape, title_text, font_size=32, bold=True, 
                            color=(255, 255, 255))
            y_pos += 1.0
        
        # 处理描述文本
        desc_elem = slide_content.find('p', class_='slide-desc')
        if desc_elem:
            desc_text = extract_text_from_element(desc_elem, lang)
            if desc_text:
                left = Inches(0.5)
                top = Inches(y_pos)
                width = Inches(9)
                height = Inches(0.6)
                desc_shape = slide.shapes.add_textbox(left, top, width, height)
                add_text_to_shape(desc_shape, desc_text, font_size=16, 
                                color=(136, 136, 136))
                y_pos += 0.8
        
        # 处理两列布局
        cols_2 = slide_content.find('div', class_='cols-2')
        if cols_2:
            # 左列内容
            left_col = cols_2.find('div')
            if left_col:
                # 处理卡片
                cards = left_col.find_all(['div'], class_=['pain-card', 'sol-card', 'card'])
                for card in cards:
                    card_text = extract_text_from_element(card, lang)
                    if card_text:
                        left = Inches(0.5)
                        top = Inches(y_pos)
                        width = Inches(5.5)
                        height = Inches(1.2)
                        card_shape = slide.shapes.add_textbox(left, top, width, height)
                        # 判断卡片类型
                        if 'pain-card' in card.get('class', []):
                            color = (255, 51, 51)  # 红色
                        elif 'sol-card' in card.get('class', []):
                            color = (204, 255, 0)  # 绿色
                        else:
                            color = (255, 255, 255)  # 白色
                        add_text_to_shape(card_shape, card_text, font_size=14, color=color)
                        y_pos += 1.3
                
                # 处理列表
                lists = left_col.find_all('ul')
                for ul in lists:
                    items = ul.find_all('li')
                    list_items = []
                    for li in items:
                        item_text = extract_text_from_element(li, lang)
                        if item_text:
                            list_items.append(f"• {item_text}")
                    if list_items:
                        list_text = "\n".join(list_items)
                        left = Inches(0.5)
                        top = Inches(y_pos)
                        width = Inches(5.5)
                        height = Inches(min(len(list_items) * 0.4, 2.5))
                        list_shape = slide.shapes.add_textbox(left, top, width, height)
                        add_text_to_shape(list_shape, list_text, font_size=14, color=(187, 187, 187))
                        y_pos += min(len(list_items) * 0.4, 2.5) + 0.2
            
            # 右列内容（指标框等）
            metric_box = cols_2.find('div', class_='metric-box')
            if metric_box:
                metric_val = metric_box.find('div', class_='metric-val')
                metric_label = metric_box.find('div', class_='metric-label')
                metric_sub = metric_box.find('div', class_='metric-sub')
                
                val_text = extract_text_from_element(metric_val, lang) if metric_val else ""
                label_text = extract_text_from_element(metric_label, lang) if metric_label else ""
                sub_text = extract_text_from_element(metric_sub, lang) if metric_sub else ""
                
                full_text = f"{val_text}\n{label_text}"
                if sub_text:
                    full_text += f"\n{sub_text}"
                
                if full_text.strip():
                    left = Inches(6.5)
                    top = Inches(2)
                    width = Inches(3)
                    height = Inches(2.5)
                    metric_shape = slide.shapes.add_textbox(left, top, width, height)
                    add_text_to_shape(metric_shape, full_text, font_size=20, bold=True, 
                                    color=(204, 255, 0), alignment=PP_ALIGN.CENTER)
        
        # 处理三列指标
        cols_3 = slide_content.find('div', class_='cols-3-metrics')
        if cols_3:
            metric_boxes = cols_3.find_all('div', class_='metric-box')
            x_start = 0.5
            width_per_box = 2.8
            top_pos = 3.5
            for i, box in enumerate(metric_boxes):
                metric_val = box.find('div', class_='metric-val')
                metric_label = box.find('div', class_='metric-label')
                metric_sub = box.find('div', class_='metric-sub')
                
                val_text = extract_text_from_element(metric_val, lang) if metric_val else ""
                label_text = extract_text_from_element(metric_label, lang) if metric_label else ""
                sub_text = extract_text_from_element(metric_sub, lang) if metric_sub else ""
                
                full_text = f"{val_text}\n{label_text}"
                if sub_text:
                    full_text += f"\n{sub_text}"
                
                if full_text.strip():
                    left = Inches(x_start + i * width_per_box)
                    top = Inches(top_pos)
                    width = Inches(2.5)
                    height = Inches(2)
                    metric_shape = slide.shapes.add_textbox(left, top, width, height)
                    add_text_to_shape(metric_shape, full_text, font_size=18, bold=True, 
                                    color=(204, 255, 0), alignment=PP_ALIGN.CENTER)
        
        # 处理网格布局
        grid_4 = slide_content.find('div', class_='grid-4')
        if grid_4:
            cards = grid_4.find_all(['div'], class_=['pain-card', 'sol-card', 'card'])
            x_positions = [0.5, 5.0]
            y_start = 2.5
            row_height = 1.2
            
            for i, card in enumerate(cards):
                card_text = extract_text_from_element(card, lang)
                if card_text:
                    # 判断是否跨列
                    span_2 = 'grid-column: span 2' in str(card.get('style', ''))
                    if span_2:
                        left = Inches(0.5)
                        width = Inches(9)
                    else:
                        col = i % 2
                        left = Inches(x_positions[col])
                        width = Inches(4.2)
                    
                    row = i // 2 if not span_2 else i
                    top = Inches(y_start + row * row_height)
                    height = Inches(row_height)
                    
                    card_shape = slide.shapes.add_textbox(left, top, width, height)
                    if 'pain-card' in card.get('class', []):
                        color = (255, 51, 51)
                    elif 'sol-card' in card.get('class', []):
                        color = (204, 255, 0)
                    else:
                        color = (255, 255, 255)
                    add_text_to_shape(card_shape, card_text, font_size=14, color=color)
        
        # 处理单独的卡片列表
        if not cols_2 and not grid_4:
            cards = slide_content.find_all(['div'], class_=['pain-card', 'sol-card', 'card'])
            for card in cards:
                if card.find_parent(['div'], class_=['cols-2', 'grid-4']):
                    continue  # 已处理
                card_text = extract_text_from_element(card, lang)
                if card_text:
                    left = Inches(0.5)
                    top = Inches(y_pos)
                    width = Inches(9)
                    height = Inches(1)
                    card_shape = slide.shapes.add_textbox(left, top, width, height)
                    if 'pain-card' in card.get('class', []):
                        color = (255, 51, 51)
                    elif 'sol-card' in card.get('class', []):
                        color = (204, 255, 0)
                    else:
                        color = (255, 255, 255)
                    add_text_to_shape(card_shape, card_text, font_size=14, color=color)
                    y_pos += 1.2
            
            # 处理单独的列表
            lists = slide_content.find_all('ul')
            for ul in lists:
                if ul.find_parent(['div'], class_=['cols-2', 'grid-4']):
                    continue  # 已处理
                items = ul.find_all('li')
                list_items = []
                for li in items:
                    item_text = extract_text_from_element(li, lang)
                    if item_text:
                        list_items.append(f"• {item_text}")
                if list_items:
                    list_text = "\n".join(list_items)
                    left = Inches(0.5)
                    top = Inches(y_pos)
                    width = Inches(9)
                    height = Inches(min(len(list_items) * 0.4, 3))
                    list_shape = slide.shapes.add_textbox(left, top, width, height)
                    add_text_to_shape(list_shape, list_text, font_size=14, color=(187, 187, 187))
                    y_pos += min(len(list_items) * 0.4, 3) + 0.3

def convert_html_to_ppt(html_file, output_file, lang='cn'):
    """将HTML文件转换为PowerPoint文件"""
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = soup.find_all('div', class_='slide')
    
    # 创建PowerPoint演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print(f"找到 {len(slides)} 个幻灯片")
    
    for i, slide_element in enumerate(slides):
        print(f"正在处理第 {i+1}/{len(slides)} 个幻灯片...")
        create_slide_from_html(prs, slide_element, lang)
    
    # 保存文件
    prs.save(output_file)
    print(f"[OK] 已保存到: {output_file}")

if __name__ == '__main__':
    html_file = r'd:\2026\0125_PlaybookS\PPT.html'
    output_file_cn = r'd:\2026\0125_PlaybookS\PPT_CN.pptx'
    output_file_en = r'd:\2026\0125_PlaybookS\PPT_EN.pptx'
    
    print("=" * 50)
    print("HTML 转 PPT 转换工具")
    print("=" * 50)
    print("\n正在转换为中文版PPT...")
    convert_html_to_ppt(html_file, output_file_cn, lang='cn')
    
    print("\n正在转换为英文版PPT...")
    convert_html_to_ppt(html_file, output_file_en, lang='en')
    
    print("\n" + "=" * 50)
    print("转换完成！")
    print("=" * 50)
