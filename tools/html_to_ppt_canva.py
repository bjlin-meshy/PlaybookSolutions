#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将HTML演示文稿转换为PowerPoint PPTX文件（Canva风格优化版）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.dml.fill import FillFormat
from bs4 import BeautifulSoup
import re

def clean_text(text):
    """清理文本，移除HTML标签和多余空白"""
    if not text:
        return ""
    # 移除HTML标签
    text = re.sub(r'<[^>]+>', '', text)
    # 清理空白字符，但保留换行
    text = re.sub(r'[ \t]+', ' ', text)
    text = text.replace('\n', ' ').strip()
    return text

def extract_text_from_element(element, lang='cn'):
    """从元素中提取指定语言的文本"""
    if not element:
        return ""
    
    # 查找指定语言的span
    lang_spans = element.find_all('span', class_=lang)
    if lang_spans:
        texts = []
        for span in lang_spans:
            text = clean_text(span.get_text())
            if text:
                texts.append(text)
        if texts:
            return '\n'.join(texts) if len(texts) > 1 else texts[0]
    
    # 如果没有找到语言特定的span，返回所有文本
    text = clean_text(element.get_text())
    return text

def add_text_to_shape(shape, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """向形状添加文本"""
    if not text:
        return
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    # 清除默认段落
    text_frame.clear()
    
    # 处理多行文本
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = line.strip()
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        if color:
            p.font.color.rgb = RGBColor(*color)
        else:
            p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(4)

def create_card_shape(slide, left, top, width, height, bg_color, border_color=None, border_width=1):
    """创建卡片形状（带背景和边框）"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    
    # 设置填充
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)
    
    # 设置边框
    if border_color:
        line = shape.line
        line.color.rgb = RGBColor(*border_color)
        line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    
    # 设置圆角
    shape.adjustments[0] = 0.05  # 圆角大小
    
    return shape

def create_left_border_card(slide, left, top, width, height, border_color, bg_gradient=False):
    """创建带左侧彩色边框的卡片"""
    # 主卡片
    card = create_card_shape(slide, left, top, width, height, 
                           (17, 17, 17), (34, 34, 34), 1)
    
    # 左侧边框条
    border_left = Inches(0.01)
    border_width = Inches(0.15)
    border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         left, top, border_width, height)
    border_fill = border_shape.fill
    border_fill.solid()
    border_fill.fore_color.rgb = RGBColor(*border_color)
    border_shape.line.fill.background()
    
    # 如果背景渐变
    if bg_gradient:
        fill = card.fill
        fill.gradient()
        fill.gradient_angle = 0.0
        fill.gradient_stops[0].color.rgb = RGBColor(255, 51, 51) if border_color == (255, 51, 51) else RGBColor(204, 255, 0)
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = RGBColor(17, 17, 17)
        fill.gradient_stops[1].position = 1.0
    
    return card

def create_slide_from_html(prs, slide_element, lang='cn'):
    """从HTML创建PowerPoint幻灯片（Canva风格）"""
    slide_content = slide_element.find('div', class_='slide-content')
    
    if not slide_content:
        return
    
    # 创建新幻灯片（使用空白布局）
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色为深灰色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 8, 8)
    
    # 添加背景光效（使用形状模拟）
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                    Inches(8), Inches(-1), 
                                    Inches(4), Inches(4))
    glow1_fill = glow1.fill
    glow1_fill.solid()
    glow1_fill.fore_color.rgb = RGBColor(12, 15, 0)  # 微弱的绿色光
    glow1.line.fill.background()
    glow1.shadow.inherit = False
    
    # 检查是否是封面页或章节标题页
    is_cover = 'align-items: center' in str(slide_content.get('style', ''))
    is_section_header = 'section-header-slide' in slide_element.get('class', [])
    
    if is_cover or is_section_header:
        # 居中布局
        # 获取标题
        title_elem = slide_content.find('h1') or slide_content.find('h2')
        title_text = ""
        if title_elem:
            title_text = extract_text_from_element(title_elem, lang)
        
        # 添加主标题
        if title_text:
            left = Inches(1)
            top = Inches(2.8)
            width = Inches(8)
            height = Inches(2)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            add_text_to_shape(title_shape, title_text, font_size=56, bold=True, 
                            color=(255, 255, 255), alignment=PP_ALIGN.CENTER)
        
        # 获取描述
        desc_elem = slide_content.find('p', class_='slide-desc')
        if not desc_elem:
            desc_elem = slide_content.find('p')
        if desc_elem:
            desc_text = extract_text_from_element(desc_elem, lang)
            if desc_text:
                left = Inches(1)
                top = Inches(5)
                width = Inches(8)
                height = Inches(1)
                desc_shape = slide.shapes.add_textbox(left, top, width, height)
                add_text_to_shape(desc_shape, desc_text, font_size=22, 
                                color=(136, 136, 136), alignment=PP_ALIGN.CENTER)
    else:
        # 常规内容页
        y_pos = 0.6
        
        # 获取section tag
        section_tag = slide_content.find('div', class_='section-tag')
        if section_tag:
            section_text = extract_text_from_element(section_tag, lang)
            if section_text:
                left = Inches(0.6)
                top = Inches(y_pos)
                width = Inches(8.8)
                height = Inches(0.35)
                tag_shape = slide.shapes.add_textbox(left, top, width, height)
                add_text_to_shape(tag_shape, section_text, font_size=11, bold=True, 
                                color=(204, 255, 0))
                y_pos += 0.45
        
        # 获取标题
        title_elem = slide_content.find('h1') or slide_content.find('h2', class_='slide-title')
        if not title_elem:
            title_elem = slide_content.find('h2')
        title_text = ""
        if title_elem:
            title_text = extract_text_from_element(title_elem, lang)
        
        if title_text:
            left = Inches(0.6)
            top = Inches(y_pos)
            width = Inches(8.8)
            height = Inches(0.9)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            add_text_to_shape(title_shape, title_text, font_size=36, bold=True, 
                            color=(255, 255, 255))
            y_pos += 1.1
        
        # 处理描述文本
        desc_elem = slide_content.find('p', class_='slide-desc')
        if desc_elem:
            desc_text = extract_text_from_element(desc_elem, lang)
            if desc_text:
                left = Inches(0.6)
                top = Inches(y_pos)
                width = Inches(8.8)
                height = Inches(0.7)
                desc_shape = slide.shapes.add_textbox(left, top, width, height)
                add_text_to_shape(desc_shape, desc_text, font_size=18, 
                                color=(136, 136, 136))
                y_pos += 0.9
        
        # 处理两列布局
        cols_2 = slide_content.find('div', class_='cols-2')
        if cols_2:
            # 左列内容
            left_col = cols_2.find('div')
            if left_col:
                # 处理卡片
                cards = left_col.find_all(['div'], class_=['pain-card', 'sol-card', 'card'])
                for card in cards:
                    card_text = extract_text_from_element(card, lang)
                    if card_text:
                        left = Inches(0.6)
                        top = Inches(y_pos)
                        width = Inches(5.2)
                        height = Inches(1.3)
                        
                        # 判断卡片类型并创建
                        if 'pain-card' in card.get('class', []):
                            card_shape = create_left_border_card(slide, left, top, width, height, 
                                                               (255, 51, 51), bg_gradient=True)
                        elif 'sol-card' in card.get('class', []):
                            card_shape = create_left_border_card(slide, left, top, width, height, 
                                                               (204, 255, 0), bg_gradient=True)
                        else:
                            card_shape = create_card_shape(slide, left, top, width, height,
                                                          (17, 17, 17), (51, 51, 51), 1)
                        
                        # 添加文本
                        text_left = Inches(0.7) if 'pain-card' in card.get('class', []) or 'sol-card' in card.get('class', []) else Inches(0.65)
                        text_top = Inches(y_pos + 0.15)
                        text_width = Inches(4.8)
                        text_height = Inches(1.0)
                        text_shape = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                        
                        color = (255, 51, 51) if 'pain-card' in card.get('class', []) else (204, 255, 0) if 'sol-card' in card.get('class', []) else (255, 255, 255)
                        add_text_to_shape(text_shape, card_text, font_size=14, color=color)
                        y_pos += 1.4
                
                # 处理列表
                lists = left_col.find_all('ul')
                for ul in lists:
                    items = ul.find_all('li')
                    list_items = []
                    for li in items:
                        item_text = extract_text_from_element(li, lang)
                        if item_text:
                            list_items.append(f"• {item_text}")
                    if list_items:
                        list_text = "\n".join(list_items)
                        left = Inches(0.6)
                        top = Inches(y_pos)
                        width = Inches(5.2)
                        height = Inches(min(len(list_items) * 0.45, 2.8))
                        list_shape = slide.shapes.add_textbox(left, top, width, height)
                        add_text_to_shape(list_shape, list_text, font_size=15, color=(187, 187, 187))
                        y_pos += min(len(list_items) * 0.45, 2.8) + 0.2
            
            # 右列内容（指标框等）
            metric_box = cols_2.find('div', class_='metric-box')
            if metric_box:
                metric_val = metric_box.find('div', class_='metric-val')
                metric_label = metric_box.find('div', class_='metric-label')
                metric_sub = metric_box.find('div', class_='metric-sub')
                
                val_text = extract_text_from_element(metric_val, lang) if metric_val else ""
                label_text = extract_text_from_element(metric_label, lang) if metric_label else ""
                sub_text = extract_text_from_element(metric_sub, lang) if metric_sub else ""
                
                full_text = f"{val_text}\n{label_text}"
                if sub_text:
                    full_text += f"\n{sub_text}"
                
                if full_text.strip():
                    # 创建指标框卡片
                    card_left = Inches(6.2)
                    card_top = Inches(2.2)
                    card_width = Inches(3.2)
                    card_height = Inches(2.8)
                    metric_card = create_card_shape(slide, card_left, card_top, card_width, card_height,
                                                   (17, 17, 17), (51, 51, 51), 1)
                    
                    # 添加文本
                    text_left = Inches(6.3)
                    text_top = Inches(2.3)
                    text_width = Inches(3.0)
                    text_height = Inches(2.6)
                    metric_shape = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                    add_text_to_shape(metric_shape, full_text, font_size=22, bold=True, 
                                    color=(204, 255, 0), alignment=PP_ALIGN.CENTER)
        
        # 处理三列指标
        cols_3 = slide_content.find('div', class_='cols-3-metrics')
        if cols_3:
            metric_boxes = cols_3.find_all('div', class_='metric-box')
            x_start = 0.6
            width_per_box = 2.7
            top_pos = 3.8
            
            for i, box in enumerate(metric_boxes):
                metric_val = box.find('div', class_='metric-val')
                metric_label = box.find('div', class_='metric-label')
                metric_sub = box.find('div', class_='metric-sub')
                
                val_text = extract_text_from_element(metric_val, lang) if metric_val else ""
                label_text = extract_text_from_element(metric_label, lang) if metric_label else ""
                sub_text = extract_text_from_element(metric_sub, lang) if metric_sub else ""
                
                full_text = f"{val_text}\n{label_text}"
                if sub_text:
                    full_text += f"\n{sub_text}"
                
                if full_text.strip():
                    # 创建指标卡片
                    card_left = Inches(x_start + i * width_per_box)
                    card_top = Inches(top_pos)
                    card_width = Inches(2.5)
                    card_height = Inches(2.2)
                    metric_card = create_card_shape(slide, card_left, card_top, card_width, card_height,
                                                   (17, 17, 17), (51, 51, 51), 1)
                    
                    # 添加文本
                    text_left = Inches(x_start + i * width_per_box + 0.1)
                    text_top = Inches(top_pos + 0.1)
                    text_width = Inches(2.3)
                    text_height = Inches(2.0)
                    metric_shape = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                    add_text_to_shape(metric_shape, full_text, font_size=20, bold=True, 
                                    color=(204, 255, 0), alignment=PP_ALIGN.CENTER)
        
        # 处理网格布局
        grid_4 = slide_content.find('div', class_='grid-4')
        if grid_4:
            cards = grid_4.find_all(['div'], class_=['pain-card', 'sol-card', 'card'])
            x_positions = [0.6, 5.2]
            y_start = 2.5
            row_height = 1.3
            
            for i, card in enumerate(cards):
                card_text = extract_text_from_element(card, lang)
                if card_text:
                    # 判断是否跨列
                    span_2 = 'grid-column: span 2' in str(card.get('style', ''))
                    if span_2:
                        left = Inches(0.6)
                        width = Inches(8.8)
                        col = 0
                    else:
                        col = i % 2
                        left = Inches(x_positions[col])
                        width = Inches(4.4)
                    
                    row = i // 2 if not span_2 else i
                    top = Inches(y_start + row * row_height)
                    height = Inches(row_height)
                    
                    # 创建卡片
                    if 'pain-card' in card.get('class', []):
                        card_shape = create_left_border_card(slide, left, top, width, height, 
                                                           (255, 51, 51), bg_gradient=True)
                        text_left = Inches(left + 0.2)
                    elif 'sol-card' in card.get('class', []):
                        card_shape = create_left_border_card(slide, left, top, width, height, 
                                                           (204, 255, 0), bg_gradient=True)
                        text_left = Inches(left + 0.2)
                    else:
                        card_shape = create_card_shape(slide, left, top, width, height,
                                                      (17, 17, 17), (51, 51, 51), 1)
                        text_left = Inches(left + 0.15)
                    
                    # 添加文本
                    text_top = Inches(y_start + row * row_height + 0.15)
                    text_width = Inches(width - 0.3)
                    text_height = Inches(row_height - 0.3)
                    text_shape = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                    
                    color = (255, 51, 51) if 'pain-card' in card.get('class', []) else (204, 255, 0) if 'sol-card' in card.get('class', []) else (255, 255, 255)
                    add_text_to_shape(text_shape, card_text, font_size=14, color=color)
        
        # 处理单独的卡片和列表
        if not cols_2 and not grid_4:
            cards = slide_content.find_all(['div'], class_=['pain-card', 'sol-card', 'card'])
            for card in cards:
                if card.find_parent(['div'], class_=['cols-2', 'grid-4']):
                    continue
                card_text = extract_text_from_element(card, lang)
                if card_text:
                    left = Inches(0.6)
                    top = Inches(y_pos)
                    width = Inches(8.8)
                    height = Inches(1.2)
                    
                    if 'pain-card' in card.get('class', []):
                        card_shape = create_left_border_card(slide, left, top, width, height, 
                                                           (255, 51, 51), bg_gradient=True)
                    elif 'sol-card' in card.get('class', []):
                        card_shape = create_left_border_card(slide, left, top, width, height, 
                                                           (204, 255, 0), bg_gradient=True)
                    else:
                        card_shape = create_card_shape(slide, left, top, width, height,
                                                      (17, 17, 17), (51, 51, 51), 1)
                    
                    text_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.15), 
                                                         Inches(8.6), Inches(0.9))
                    color = (255, 51, 51) if 'pain-card' in card.get('class', []) else (204, 255, 0) if 'sol-card' in card.get('class', []) else (255, 255, 255)
                    add_text_to_shape(text_shape, card_text, font_size=14, color=color)
                    y_pos += 1.3
            
            lists = slide_content.find_all('ul')
            for ul in lists:
                if ul.find_parent(['div'], class_=['cols-2', 'grid-4']):
                    continue
                items = ul.find_all('li')
                list_items = []
                for li in items:
                    item_text = extract_text_from_element(li, lang)
                    if item_text:
                        list_items.append(f"• {item_text}")
                if list_items:
                    list_text = "\n".join(list_items)
                    left = Inches(0.6)
                    top = Inches(y_pos)
                    width = Inches(8.8)
                    height = Inches(min(len(list_items) * 0.45, 3))
                    list_shape = slide.shapes.add_textbox(left, top, width, height)
                    add_text_to_shape(list_shape, list_text, font_size=15, color=(187, 187, 187))
                    y_pos += min(len(list_items) * 0.45, 3) + 0.3

def convert_html_to_ppt(html_file, output_file, lang='cn'):
    """将HTML文件转换为PowerPoint文件"""
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = soup.find_all('div', class_='slide')
    
    # 创建PowerPoint演示文稿 - 使用16:9比例（1280x720）
    prs = Presentation()
    prs.slide_width = Emu(12800000)  # 1280px = 12800000 EMU
    prs.slide_height = Emu(7200000)   # 720px = 7200000 EMU
    
    print(f"找到 {len(slides)} 个幻灯片")
    
    for i, slide_element in enumerate(slides):
        print(f"正在处理第 {i+1}/{len(slides)} 个幻灯片...")
        create_slide_from_html(prs, slide_element, lang)
    
    # 保存文件
    prs.save(output_file)
    print(f"[OK] 已保存到: {output_file}")

if __name__ == '__main__':
    html_file = r'd:\2026\0125_PlaybookS\PPT.html'
    output_file_cn = r'd:\2026\0125_PlaybookS\PPT_CN_Canva.pptx'
    output_file_en = r'd:\2026\0125_PlaybookS\PPT_EN_Canva.pptx'
    
    print("=" * 50)
    print("HTML 转 PPT 转换工具 (Canva风格优化版)")
    print("=" * 50)
    print("\n正在转换为中文版PPT...")
    convert_html_to_ppt(html_file, output_file_cn, lang='cn')
    
    print("\n正在转换为英文版PPT...")
    convert_html_to_ppt(html_file, output_file_en, lang='en')
    
    print("\n" + "=" * 50)
    print("转换完成！")
    print("=" * 50)
